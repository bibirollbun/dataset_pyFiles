import os, pathlib, tempfile, atexit
from kaggle_secrets import UserSecretsClient
from google.cloud import bigquery

u = UserSecretsClient()

def need(name: str) -> str:
    v = u.get_secret(name)
    if not v:
        raise RuntimeError(f"Missing secret: {name}")
    print(f"✓ {name}")
    return v

# --- secrets / config ---
sa_json    = need("GCP_SERVICE_ACCOUNT_JSON")
PROJECT_ID = need("PROJECT_ID")
REGION     = need("BQ_LOCATION")   
DATASET    = need("DATASET_ID")
CONNECTION = u.get_secret("BQ_VERTEX_CONNECTION") or f"projects/{PROJECT_ID}/locations/{REGION}/connections/vertex_conn"
os.environ["BQ_VERTEX_CONNECTION"] = CONNECTION
os.environ["GOOGLE_CLOUD_PROJECT"] = PROJECT_ID

# --- temp key ---
temp_dir = "/kaggle/temp" if os.path.isdir("/kaggle/temp") else tempfile.gettempdir()
pathlib.Path(temp_dir).mkdir(parents=True, exist_ok=True)
sa_path = os.path.join(temp_dir, "gcp_sa.json")
pathlib.Path(sa_path).write_text(sa_json)
os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = sa_path
atexit.register(lambda: (os.path.exists(sa_path) and os.remove(sa_path)))

# --- BigQuery client + sanity check ---
client = bigquery.Client(project=PROJECT_ID)

# location goes here (on query), not in QueryJobConfig
rows = list(client.query("SELECT 1 AS ok", location=REGION).result())
print(rows)
print("PROJECT_ID:", PROJECT_ID, "| REGION:", REGION, "| DATASET:", DATASET)


from google.cloud import bigquery
client = bigquery.Client(project=PROJECT_ID, location=REGION)
print(list(client.query("SELECT 1 AS ok").result()))


client.query(f"""
CREATE SCHEMA IF NOT EXISTS `{PROJECT_ID}.{DATASET}`
OPTIONS(location='{REGION}');
""").result()

client.query(f"""
CREATE OR REPLACE MODEL `{PROJECT_ID}.{DATASET}.gemini_remote`
REMOTE WITH CONNECTION `{os.environ['BQ_VERTEX_CONNECTION']}`
OPTIONS(ENDPOINT = 'gemini-2.0-flash-001');
""").result()

client.query(f"""
CREATE OR REPLACE MODEL `{PROJECT_ID}.{DATASET}.embed_remote`
REMOTE WITH CONNECTION `{os.environ['BQ_VERTEX_CONNECTION']}`
OPTIONS(ENDPOINT = 'text-embedding-004');
""").result()

print("✓ remote models ready")


!pip -q install python-pptx

import os
from pptx import Presentation

SAMPLE_ROOT = "/kaggle/input/testdata"
assert os.path.exists(SAMPLE_ROOT), "Attach your sample docs dataset in Settings → Data."

rows = []
def add_unit(path, text, slide=None, page=None, panel_label=None, row_header=None, col_header=None):
    rows.append({
        "source_path": path, "slide": slide, "page": page,
        "panel_label": panel_label, "row_header": row_header, "col_header": col_header,
        "text_full": text
    })

for root, _, files in os.walk(SAMPLE_ROOT):
    for f in files:
        p = os.path.join(root, f)
        fl = f.lower()
        if fl.endswith((".txt",".md")):
            with open(p, "r", encoding="utf-8", errors="ignore") as fh:
                buf=[]
                for line in fh:
                    line=line.strip()
                    if line: buf.append(line)
                    elif buf: add_unit(p, " ".join(buf)); buf=[]
                if buf: add_unit(p, " ".join(buf))
        elif fl.endswith(".pptx"):
            try:
                prs = Presentation(p)
                for i, slide in enumerate(prs.slides, start=1):
                    texts=[getattr(s,"text","") for s in slide.shapes if hasattr(s,"text") and s.text]
                    txt="\n".join(t.strip() for t in texts if t and t.strip())
                    if txt: add_unit(p, txt, slide=i)
            except Exception as e:
                print("PPTX parse failed:", p, "→", e)

units_table = f"{PROJECT_ID}.{DATASET}.units_demo"
client.query(f"""
CREATE OR REPLACE TABLE `{units_table}` (
  source_path STRING, slide INT64, page INT64,
  panel_label STRING, row_header STRING, col_header STRING,
  text_full STRING
)""").result()

load = client.load_table_from_json(rows, units_table)
load.result()
print("✓ loaded rows:", len(rows))


# ==== Build embeddings table from units_demo ====
import os, json, pathlib

from google.cloud import bigquery



# BigQuery client
client = bigquery.Client(project=PROJECT_ID, location=REGION)

# Tables / models
units_table = f"{PROJECT_ID}.{DATASET}.units_demo"
vec_table   = f"{PROJECT_ID}.{DATASET}.units_demo_vec"
embed_model = f"{PROJECT_ID}.{DATASET}.embed_remote"   
# Build embeddings
client.query(f"""
CREATE OR REPLACE TABLE `{vec_table}` AS
SELECT
  GENERATE_UUID() AS unit_id,
  source_path, slide, page, panel_label, row_header, col_header, text_full,
  (
    SELECT e.ml_generate_embedding_result
    FROM ML.GENERATE_EMBEDDING(
      MODEL `{embed_model}`,
      (SELECT text_full AS content)
    ) AS e
  ) AS embedding
FROM `{units_table}`
WHERE text_full IS NOT NULL AND LENGTH(text_full) > 0
""").result()

print("✓ built embeddings table:", vec_table)


question = "Summarize the key points and any recommended actions mentioned in these docs."

# Query embedding 
q_emb = list(client.query(f"""
DECLARE q STRING DEFAULT @q;
SELECT ml_generate_embedding_result
FROM ML.GENERATE_EMBEDDING(
  MODEL `{embed_model}`,
  (SELECT q AS content)
)""", job_config=bigquery.QueryJobConfig(
    query_parameters=[bigquery.ScalarQueryParameter("q","STRING", question)]
)).result())[0].ml_generate_embedding_result

# Dot-product similarity
topk = list(client.query(f"""
WITH q AS (SELECT {q_emb} AS qvec)
SELECT unit_id, source_path, slide, page, text_full,
       (SELECT SUM(a*b)
          FROM UNNEST(v.embedding) a WITH OFFSET i
          JOIN UNNEST(q.qvec)     b WITH OFFSET j
          ON i=j) AS score
FROM `{vec_table}` v, q
ORDER BY score DESC
LIMIT 5
""").result())

print("TopK =", len(topk))
for i, r in enumerate(topk, 1):
    cite = f"(slide {r.slide})" if r.slide else (f"(page {r.page})" if r.page else "")
    print(f"{i}. {round(r.score,3)} {cite} — {os.path.basename(r.source_path)}")


gen_model = f"{PROJECT_ID}.{DATASET}.gemini_remote"

ctx = []
for r in topk:
    cite = f"(slide {r.slide})" if r.slide else (f"(page {r.page})" if r.page else "")
    ctx.append(f"[{r.unit_id}] {cite} {r.source_path}: {r.text_full}")

prompt = f"""You are an enterprise RAG assistant. Answer using ONLY the context below.
Cite evidence with [unit_id] and include slide/page when present.
Question: {question}
Context:
{chr(10).join(ctx)}
"""

ans = list(client.query(f"""
DECLARE p STRING DEFAULT @p;
SELECT *
FROM ML.GENERATE_TEXT(
  MODEL `{gen_model}`,
  (SELECT p AS prompt),
  STRUCT(0.2 AS temperature, 512 AS max_output_tokens)
)
""", job_config=bigquery.QueryJobConfig(
    query_parameters=[bigquery.ScalarQueryParameter("p","STRING", prompt)]
)).result())

print("\n=== ANSWER ===\n")
print(ans[0].ml_generate_text_result)  


# === Persist RAG Q/A + context and zip as Kaggle submission artifact ===
import os, json, pathlib, zipfile, datetime

ART_DIR = pathlib.Path("/kaggle/working/artifacts")
ART_DIR.mkdir(parents=True, exist_ok=True)

# 1) Recreate the exact SQL from earlier
GEN_SQL = f"""
DECLARE p STRING DEFAULT @p;
SELECT *
FROM ML.GENERATE_TEXT(
  MODEL `{gen_model}`,
  (SELECT p AS prompt),
  STRUCT(0.2 AS temperature, 512 AS max_output_tokens)
)
"""

# 2) Extract answer text robustly from BigQuery result
res = ans[0].ml_generate_text_result
answer_text = None
if isinstance(res, dict):
    # common BQ AI shapes
    answer_text = (
        (res.get("predictions") or [{}])[0].get("content") or
        (res.get("results") or [{}])[0].get("output_text") or
        res.get("content") or
        res.get("output_text")
    )
if not answer_text:
    # fallback to string
    answer_text = str(res)

# 3) Flatten the context used
ctx_lines = []
for r in topk:
    cite = f"(slide {getattr(r, 'slide', None)})" if getattr(r, 'slide', None) else (f"(page {getattr(r, 'page', None)})" if getattr(r, 'page', None) else "")
    unit = getattr(r, 'unit_id', None) or "<no_unit_id>"
    src  = getattr(r, 'source_path', None) or "<no_source>"
    text = getattr(r, 'text_full', None) or ""
    ctx_lines.append(f"[{unit}] {cite} {src}: {text}")

# 4) Write artifacts (plain text for easy viewing)
(ART_DIR / "question.txt").write_text(str(question))
(ART_DIR / "prompt.txt").write_text(prompt)
(ART_DIR / "context.txt").write_text("\n".join(ctx_lines))
(ART_DIR / "answer.txt").write_text(answer_text)
(ART_DIR / "generate_text.sql").write_text(GEN_SQL)

# Minimal run metadata (no secrets)
meta = {
    "generated_at_utc": datetime.datetime.utcnow().isoformat() + "Z",
    "project_id": os.getenv("GOOGLE_CLOUD_PROJECT", "<unset>"),
    "region": os.getenv("BQ_LOCATION", "<unset>"),
    "dataset": os.getenv("DATASET_ID", "<unset>"),
    "model": gen_model,
    "rows_in_context": len(ctx_lines)
}
(ART_DIR / "run_metadata.json").write_text(json.dumps(meta, indent=2))

# 5) Zip everything to a single submission file
zip_path = "/kaggle/working/bqai_rag_submission.zip"
with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as z:
    for p in ART_DIR.iterdir():
        if p.is_file():
            z.write(p, arcname=f"artifacts/{p.name}")

print("Created submission file:", zip_path)
print("Saved:", [p.name for p in ART_DIR.iterdir() if p.is_file()])


