# This Python 3 environment comes with many helpful analytics libraries installed
# It is defined by the kaggle/python Docker image: https://github.com/kaggle/docker-python
# For example, here's several helpful packages to load

import numpy as np # linear algebra
import pandas as pd # data processing, CSV file I/O (e.g. pd.read_csv)

# Input data files are available in the read-only "../input/" directory
# For example, running this (by clicking run or pressing Shift+Enter) will list all files under the input directory

import os
for dirname, _, filenames in os.walk('/kaggle/input'):
    for filename in filenames:
        print(os.path.join(dirname, filename))

# You can write up to 20GB to the current directory (/kaggle/working/) that gets preserved as output when you create a version using "Save & Run All" 
# You can also write temporary files to /kaggle/temp/, but they won't be saved outside of the current session


# -----------------------------
# AI STUDY HELPER (FINAL) - Kaggle notebook single cell
# Features: question auto-clear, PDF/DOCX/PPTX/Image uploads, combine file-text+question,
# produce multiple answer styles (Simple, Detailed, Exam, Study Notes) using Google Gemini
# -----------------------------

# --- auto-install helper (best-effort) ---
import sys, subprocess, importlib
def safe_import(pkg, name=None):
    nm = name or pkg
    try:
        return importlib.import_module(nm)
    except Exception:
        try:
            print(f"Installing {pkg} ...")
            subprocess.check_call([sys.executable, "-m", "pip", "install", pkg], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            return importlib.import_module(nm)
        except Exception as e:
            print(f"Could not install/import {pkg}: {e}")
            return None

# --- try required libs ---
ipywidgets = safe_import("ipywidgets", "ipywidgets")
from IPython.display import display, clear_output
widgets = ipywidgets

PyPDF2 = safe_import("PyPDF2", "PyPDF2")
docx = safe_import("python-docx", "docx")
pptx = safe_import("python-pptx", "pptx")
PIL = safe_import("Pillow", "PIL")
pytesseract = safe_import("pytesseract", "pytesseract")  # OCR (may not be present in Kaggle)

# Gemini client
genai_pkg = safe_import("google-genai", "google")
try:
    from google import genai
    from google.genai import types
except Exception:
    genai = None

# Kaggle secrets
kaggle_secrets_pkg = safe_import("kaggle-secrets", "kaggle_secrets")
API_KEY = None
try:
    from kaggle_secrets import UserSecretsClient
    user_secrets = UserSecretsClient()
    API_KEY = user_secrets.get_secret("GOOGLE_API_KEY")
except Exception:
    API_KEY = None

# Inform if missing
if genai is None or API_KEY is None:
    print("⚠️ Note: google-genai package or GOOGLE_API_KEY missing. Set 'GOOGLE_API_KEY' in Kaggle Secrets and ensure google-genai installed.")
else:
    client = genai.Client(api_key=API_KEY)

# -----------------------------
# Helper: file text extraction
# -----------------------------

import io

def extract_text_from_pdf_bytes(b):
    if PyPDF2 is None:
        return "(PDF parser not available: PyPDF2 missing.)"
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(b))
        pages = []
        for p in reader.pages:
            try:
                t = p.extract_text() or ""
                pages.append(t)
            except Exception:
                pages.append("")
        return "\n".join(pages).strip()
    except Exception as e:
        return f"(PDF extraction error: {e})"

def extract_text_from_docx_bytes(b):
    if docx is None:
        return "(DOCX parser not available: python-docx missing.)"
    try:
        from docx import Document
        doc = Document(io.BytesIO(b))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        return f"(DOCX extraction error: {e})"

def extract_text_from_pptx_bytes(b):
    if pptx is None:
        return "(PPTX parser not available: python-pptx missing.)"
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(b))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
                except Exception:
                    continue
        return "\n".join(texts)
    except Exception as e:
        return f"(PPTX extraction error: {e})"

def extract_text_from_image_bytes(b):
    if PIL is None:
        return "(Pillow not available: can't open image.)"
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(b)).convert("RGB")
        if pytesseract is None:
            return "(pytesseract not installed: OCR not available.)"
        text = pytesseract.image_to_string(img)
        return text
    except Exception as e:
        return f"(Image OCR error: {e})"

def extract_text_from_upload_item(item):
    """
    item may be one of:
      - dict-like from ipywidgets.FileUpload: {'metadata':..., 'content': b'...'}
      - or older formats: tuple/list; adapt below
    """
    try:
        # If dict with 'content' and 'metadata'
        if isinstance(item, dict) and 'content' in item:
            name = item['metadata'].get('name') if item.get('metadata') else item.get('name')
            content = item['content']
        # If tuple/list like (name, {..})
        elif isinstance(item, (list, tuple)) and len(item) >= 2:
            name = item[0]
            content = item[1].get('content') if isinstance(item[1], dict) else item[1]
        else:
            # fallback: if item itself is bytes or has 'name'/'content' attributes
            name = getattr(item, 'name', 'uploaded_file')
            content = getattr(item, 'content', None) or getattr(item, 'read', lambda: b'')()
        name = name.lower()
        # call appropriate parser
        if name.endswith(".pdf"):
            return name, extract_text_from_pdf_bytes(content)
        if name.endswith(".docx"):
            return name, extract_text_from_docx_bytes(content)
        if name.endswith(".pptx"):
            return name, extract_text_from_pptx_bytes(content)
        if any(name.endswith(ext) for ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]):
            return name, extract_text_from_image_bytes(content)
        # fallback: try decode
        try:
            return name, content.decode("utf-8", errors="ignore")
        except Exception:
            return name, "(Unable to extract text from this file type.)"
    except Exception as e:
        return "unknown", f"(Upload processing error: {e})"

# -----------------------------
# Build Gemini prompts (multiple styles)
# -----------------------------
def build_prompts(base_text, question, lang="English"):
    # base_text: combined file text (string)
    # question: user question (string)
    context = ""
    if base_text:
        context = f"Context from uploaded files:\n{base_text}\n\n"
    q = question or ""
    # create variants
    simple = f"Language: {lang}\nAnswer briefly for a student:\n\n{context}Question: {q}\n\nShort answer:"
    detailed = f"Language: {lang}\nProvide a clear, detailed explanation with examples, steps and a 3-sentence summary at the start:\n\n{context}Question: {q}\n\nDetailed explanation:"
    exam = f"Language: {lang}\nProvide exam-style responses: 1) short answer (2-3 lines) 2) long answer (6-8 lines) 3) extra important points:\n\n{context}Question: {q}\n\nExam-style answer:"
    notes = f"Language: {lang}\nCreate study notes with headings, bullet points and a 2-line summary at end:\n\n{context}Question: {q}\n\nStudy notes:"
    return {"Simple": simple, "Detailed": detailed, "Exam": exam, "StudyNotes": notes}

def ask_gemini(prompt, model="gemini-2.0-flash"):
    if genai is None or API_KEY is None:
        return "(Gemini client or API key not available. Please set GOOGLE_API_KEY in Kaggle Secrets.)"
    try:
        resp = client.models.generate_content(model=model, contents=prompt)
        return getattr(resp, "text", str(resp))
    except Exception as e:
        return f"(Gemini API error: {e})"

# -----------------------------
# UI: ipywidgets
# -----------------------------
from ipywidgets import Text, Dropdown, Button, FileUpload, Output, VBox, HBox, Label, Checkbox, SelectMultiple

question_box = Text(placeholder="Type your study question (or leave blank to use only file content)", layout=widgets.Layout(width="70%"))
style_dropdown = Dropdown(options=["All","Simple","Detailed","Exam","StudyNotes"], value="All", description="Output style:")
ask_btn = Button(description="Ask AI", button_style="success")
clear_btn = Button(description="Clear box", button_style="warning")
file_upload = FileUpload(accept=".pdf,.docx,.pptx,.png,.jpg,.jpeg,.txt", multiple=True)
process_btn = Button(description="Process Uploads", button_style="info")
output = Output(layout={'border':'1px solid lightgray'})

# internal storage
uploaded_texts = []  # list of (filename, text)

# Handlers
def on_process_uploads(b):
    global uploaded_texts
    uploaded_texts = []
    with output:
        clear_output()
        # adapt to returned format
        val = file_upload.value
        if not val:
            print("No files uploaded.")
            return
        # val may be dict-like or tuple/list - try both
        items = []
        try:
            # dict-like with filenames as keys
            if isinstance(val, dict):
                for k,v in val.items():
                    # v may be dict with 'metadata' and 'content'
                    items.append(v)
            elif isinstance(val, (list, tuple)):
                items = list(val)
            else:
                # unknown structure: try iterate
                for it in val:
                    items.append(it)
        except Exception:
            items = list(val) if hasattr(val, "__iter__") else []
        # process items
        for it in items:
            name, txt = extract_text_from_upload_item(it)
            uploaded_texts.append((name, txt))
            print(f"Processed: {name} (text length: {len(txt)} chars)")
    # keep existing question text, do not overwrite

def on_ask(b):
    with output:
        clear_output()
        # prepare combined text
        combined = ""
        for name, txt in uploaded_texts:
            combined += f"\n--- FILE: {name} ---\n{txt}\n"
        q = question_box.value.strip()
        if not q and not combined.strip():
            print("Please type a question or upload files then click Process Uploads.")
            return
        # choose styles
        style_choice = style_dropdown.value
        prompts = build_prompts(combined, q, lang="English")
        styles_to_run = ["Simple","Detailed","Exam","StudyNotes"] if style_choice=="All" else [style_choice if style_choice!="StudyNotes" else "StudyNotes"]
        # ask Gemini for each requested style
        print("AI is thinking... this may take a few seconds.")
        for s in styles_to_run:
            prompt = prompts["StudyNotes"] if s=="StudyNotes" else prompts[s]
            result = ask_gemini(prompt)
            print("\n" + ("="*8) + f" {s} " + ("="*8) + "\n")
            print(result[:20000])  # print safely (limit)
        # clear question box automatically
        question_box.value = ""
        print("\nDone. Question box cleared. Ask another question or upload more files.")

def on_clear(b):
    question_box.value = ""
    with output:
        clear_output()
        print("Question box cleared.")

# wire handlers
process_btn.on_click(on_process_uploads)
ask_btn.on_click(on_ask)
clear_btn.on_click(on_clear)

# layout & display
display(VBox([
    Label("AI Study Helper — Type a question and/or upload files (PDF/DOCX/PPTX/IMAGE)."),
    HBox([question_box, style_dropdown, ask_btn, clear_btn]),
    HBox([Label("Upload files (multiple):"), file_upload, process_btn]),
    output
]))

