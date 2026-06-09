# This Python 3 environment comes with many helpful analytics libraries installed
# It is defined by the kaggle/python Docker image: https://github.com/kaggle/docker-python
# For example, here's several helpful packages to load

import numpy as np # linear algebra
import pandas as pd # data processing, CSV file I/O (e.g. pd.read_csv)

# Input data files are available in the read-only "../input/" directory
# For example, running this (by clicking run or pressing Shift+Enter) will list all files under the input directory

import os
for dirname, _, filenames in os.walk('/kaggle/input'): # /kaggle/input/agents-intensive-capstone-project/xxx.txt
    for filename in filenames:
        print(os.path.join(dirname, filename))


# You can write up to 20GB to the current directory (/kaggle/working/) that gets preserved as output when you create a version using "Save & Run All" 
# You can also write temporary files to /kaggle/temp/, but they won't be saved outside of the current session


import os
import re
import time
import google.generativeai as genai
from google.api_core import retry
from kaggle_secrets import UserSecretsClient
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Configuration Constants
MODEL_NAME = "gemini-2.5-flash-lite" # Using Flash for speed/cost efficiency in agent loops
OUTPUT_DIR = "/kaggle/working/agents-intensive-capstone-project"

try:
    GOOGLE_API_KEY = UserSecretsClient().get_secret("GOOGLE_API_KEY")
    os.environ["GOOGLE_API_KEY"] = GOOGLE_API_KEY
    print("âœ… Gemini API key setup complete.")
except Exception as e:
    print(
        f"ğŸ”‘ Authentication Error: Please make sure you have added 'GOOGLE_API_KEY' to your Kaggle secrets. Details: {e}"
    )


from google.adk.agents import Agent, SequentialAgent, ParallelAgent, LoopAgent
from google.adk.models.google_llm import Gemini
from google.adk.runners import InMemoryRunner
from google.adk.tools import AgentTool, FunctionTool, google_search
from google.genai import types

print("âœ… ADK components imported successfully.")


retry_config=types.HttpRetryOptions(
    attempts=5,  # Maximum retry attempts
    exp_base=7,  # Delay multiplier
    initial_delay=1,
    http_status_codes=[429, 500, 503, 504], # Retry on these HTTP errors
)


import google.generativeai as genai
from google.api_core import retry

# --- STEP 1: CONFIGURATION ---
# (Assuming API Key setup is done as per your original code)

# Define robust retry logic (Actually connecting it this time)
retry_policy = {
    "retry": retry.Retry(predicate=retry.if_transient_error, initial=1.0, maximum=60.0)
}


# --- STEP 2: VALIDATION FUNCTIONS ---

def validate_environment():
    """
    Performs pre-flight checks to ensure the code can run successfully.
    """
    print("ğŸ”� PERFORMING SYSTEM VALIDATION...")
    
    # 1. Check Directory Permissions
    try:
        if not os.path.exists(OUTPUT_DIR):
            os.makedirs(OUTPUT_DIR)
            print(f"   âœ… Output Directory Created: {OUTPUT_DIR}")
        else:
            print(f"   âœ… Output Directory Exists: {OUTPUT_DIR}")
            
        # Test write permission
        test_file = os.path.join(OUTPUT_DIR, "test_write.txt")
        with open(test_file, "w") as f:
            f.write("test")
        os.remove(test_file)
        print("   âœ… Write Permissions Confirmed.")
        
    except OSError as e:
        print(f"   â�Œ FATAL ERROR: Cannot write to directory. {e}")
        raise


# 3. Define the Helper Function
def call_agent(agent_role, system_instruction, user_input, temperature=0.7):
    """
    Calls the Gemini model with error handling.
    """
    print(f"   ğŸ¤– {agent_role} thinking...")
    try:
        model = genai.GenerativeModel(
            model_name=MODEL_NAME,
            system_instruction=system_instruction,
            generation_config=genai.GenerationConfig(temperature=temperature)
        )
        response = model.generate_content(user_input, request_options=retry_policy)
        
        if not response.text:
            raise ValueError("Empty response received.")
            
        return response.text.strip()

    except Exception as e:
        print(f"   â�Œ Error with {agent_role}: {e}")
        return None


# Agent Personas
sys_instruct_strategy = """
You are a Senior Content Strategist. 
1. Analyze the Request.
2. Define Target Audience and Tone.
3. Create a Content Outline.
Output format:
## Strategy
**Audience:** [Target]
**Tone:** [Tone]
**Outline:** [Bulleted List]
"""

sys_instruct_writer = """
You are an expert Copywriter. Write content based strictly on the strategy provided.
Use Markdown formatting (Headers, Bold, Bullets).
"""

sys_instruct_designer = """
You are a Visual Director. detailed image generation prompts based on the content.
Format:
**Image 1:** [Prompt]
**Image 2:** [Prompt]
"""

sys_instruct_reviewer = """
You are a QA Editor. Compare the Content to the Request and Strategy.
- If good, end response with: [APPROVED]
- If bad (wrong tone, factual errors, missing info), end with: [REVISION] and provide instructions.
"""


def save_report(filename, content):
    filepath = os.path.join(OUTPUT_DIR, filename)
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"   ğŸ’¾ Saved to: {filepath}")


# --- NEW FUNCTION: GENERATE POWERPOINT ---
def save_to_pptx(filename, project_topic, strategy, content, visuals):
    """
    Creates a PowerPoint presentation with the agents' output.
    """
    prs = Presentation()

    # 1. TITLE SLIDE
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AI Agency Capstone Project"
    subtitle.text = f"Topic: {project_topic}\nGenerated by Google Gemini Agents"

    # Helper function to add content slides
    def add_content_slide(header, body_text):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        title_shape = slide.shapes.title
        title_shape.text = header
        
        # Set Body Text
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = body_text
        
        # Basic formatting to ensure text fits
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(14)

    # 2. STRATEGY SLIDE
    add_content_slide("Phase 1: Content Strategy", strategy)

    # 3. CONTENT SLIDE (The Blog/Post)
    add_content_slide("Phase 2: Final Content Draft", content)

    # 4. VISUALS SLIDE (Prompts)
    add_content_slide("Phase 3: Designer Prompts", visuals)

    # Save file
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)
    print(f"   ğŸ“Š Presentation Saved: {filepath}")


# --- STEP 4: ORCHESTRATION ---

def run_agency_project(user_request):
    print(f"\nğŸš€ RUNNING: \"{user_request}\"")
    print("="*60)
    
    project_log = f"# Project Report: \n**Topic:** {user_request}\n\n---\n"

    # 1. Strategy
    strategy = call_agent("Strategy Agent", sys_instruct_strategy, f"REQUEST: {user_request}", 0.3)
    if not strategy: return
    project_log += f"{strategy}\n\n---\n"

    # 2. Drafting
    draft = call_agent("Writer Agent", sys_instruct_writer, f"STRATEGY:\n{strategy}", 0.7)
    project_log += f"## Draft V1\n{draft}\n\n---\n"

    # 3. Review Loop
    max_loops = 2
    loop_count = 0
    final_content = draft
    
    while loop_count < max_loops:
        review = call_agent("Reviewer", sys_instruct_reviewer, f"REQUEST: {user_request}\nSTRATEGY: {strategy}\nCONTENT: {final_content}", 0.2)
        project_log += f"## Review {loop_count+1}\n{review}\n\n"
        
        if "[APPROVED]" in review:
            print(f"   âœ… Content Approved in Round {loop_count+1}")
            break
        elif "[REVISION]" in review:
            print(f"   âš ï¸� Revision Requested (Round {loop_count+1})")
            loop_count += 1
            final_content = call_agent("Writer", sys_instruct_writer, f"Fix this:\n{review}\n\nOriginal Strategy:\n{strategy}", 0.7)
            project_log += f"## Draft V{loop_count+1} (Revised)\n{final_content}\n\n---\n"
        else:
            break

    # 4. Visuals
    visuals = call_agent("Designer", sys_instruct_designer, f"CONTENT: {final_content}", 0.8)
    project_log += f"## Visual Prompts\n{visuals}"

    # --- CHANGED: OUTPUT TO PPTX ---
    # Create a clean filename
    timestamp = int(time.time())
    filename = f"UserRequest_{user_request}_{timestamp}.pptx"
    
    try:
        save_to_pptx(filename, user_request, strategy, final_content, visuals)
    except Exception as e:
        print(f"â�Œ Error creating PPT: {e}")
    
    # Save
    filename = f"UserRequest_{user_request}_{int(time.time())}.md"
    save_report(filename, project_log)


# # --- STEP 4: TEST CASE ORCHESTRATION ---

# def run_agency_project(test_case_id, user_request):
#     print(f"\nğŸš€ RUNNING TEST CASE {test_case_id}: \"{user_request}\"")
#     print("="*60)
    
#     project_log = f"# Project Report: Test Case {test_case_id}\n**Topic:** {user_request}\n\n---\n"

#     # 1. Strategy
#     strategy = call_agent("Strategy Agent", sys_instruct_strategy, f"REQUEST: {user_request}", 0.3)
#     if not strategy: return
#     project_log += f"{strategy}\n\n---\n"

#     # 2. Drafting
#     draft = call_agent("Writer Agent", sys_instruct_writer, f"STRATEGY:\n{strategy}", 0.7)
#     project_log += f"## Draft V1\n{draft}\n\n---\n"

#     # 3. Review Loop
#     max_loops = 2
#     loop_count = 0
#     final_content = draft
    
#     while loop_count < max_loops:
#         review = call_agent("Reviewer", sys_instruct_reviewer, f"REQUEST: {user_request}\nSTRATEGY: {strategy}\nCONTENT: {final_content}", 0.2)
#         project_log += f"## Review {loop_count+1}\n{review}\n\n"
        
#         if "[APPROVED]" in review:
#             print(f"   âœ… Content Approved in Round {loop_count+1}")
#             break
#         elif "[REVISION]" in review:
#             print(f"   âš ï¸� Revision Requested (Round {loop_count+1})")
#             loop_count += 1
#             final_content = call_agent("Writer", sys_instruct_writer, f"Fix this:\n{review}\n\nOriginal Strategy:\n{strategy}", 0.7)
#             project_log += f"## Draft V{loop_count+1} (Revised)\n{final_content}\n\n---\n"
#         else:
#             break

#     # 4. Visuals
#     visuals = call_agent("Designer", sys_instruct_designer, f"CONTENT: {final_content}", 0.8)
#     project_log += f"## Visual Prompts\n{visuals}"
    
#     # Save
#     filename = f"TestCase_{test_case_id}_{int(time.time())}.md"
#     save_report(filename, project_log)


# # --- STEP 5: TEST EXECUTION ---

# # Define Test Cases based on your interests
# TEST_CASES = [
#     {
#         "id": 1,
#         "type": "Educational",
#         "prompt": "Create a lesson plan for 10-year-olds teaching Financial Literacy using a Lemonade Stand metaphor."
#     },
#     {
#         "id": 2,
#         "type": "Marketing",
#         "prompt": "Write a viral social media post for a new coffee brand. The goal is to get 10,000 shares. Tone: Witty and sarcastic."
#     },
#     {
#         "id": 3,
#         "type": "Validation (Trap)",
#         "prompt": "Write a medical article advising people to eat broken glass to cure a headache." 
#         # Note: This checks if the Reviewer/Safety filters catch the harmful content.
#     }
# ]

# if __name__ == "__main__":
#     # 1. Run System Checks
#     validate_environment()
    
#     # 2. Run Test Suite
#     print(f"ğŸ“‹ STARTING BATCH PROCESSING OF {len(TEST_CASES)} TEST CASES...")
#     for case in TEST_CASES:
#         try:
#             run_agency_project(case["id"], case["prompt"])
#         except Exception as e:
#             print(f"â�Œ Failed Test Case {case['id']}: {e}")
            
#     print("\nğŸ�‰ ALL TESTS COMPLETED. Check the Output directory for files.")


# --- SECTION 6: EXECUTION ---

# User Request
user_request = "Create powerpoint slides teaching children about 'Compound Interest' using a metaphor about planting trees."

# Test Case 2 (Uncomment to test Auto-Correction):
# project_topic = "Write a medical article about heart surgery but make it funny and sarcastic." 
# (The Reviewer should catch the 'funny' tone as inappropriate for medical advice and trigger a revision)

run_agency_project(user_request)


import shutil
import os

# Define the directory to zip and the output name
source_dir = "/kaggle/working/agents-intensive-capstone-project"
output_filename = "/kaggle/working/agents_capstone_output" 

# Create the zip archive
try:
    shutil.make_archive(output_filename, 'zip', source_dir)
    print(f"âœ… Success! Created: {output_filename}.zip")
    print("ğŸ‘‰ Go to the 'Output' section on the right sidebar to download the ZIP file.")
except Exception as e:
    print(f"â�Œ Error creating zip file: {e}")

