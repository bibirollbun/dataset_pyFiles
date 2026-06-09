# This Python 3 environment comes with many helpful analytics libraries installed
# It is defined by the kaggle/python Docker image: https://github.com/kaggle/docker-python
# For example, here's several helpful packages to load

import numpy as np # linear algebra
import pandas as pd # data processing, CSV file I/O (e.g. pd.read_csv)

# Input data files are available in the read-only "../input/" directory
# For example, running this (by clicking run or pressing Shift+Enter) will list all files under the input directory

import os
for dirname, _, filenames in os.walk('/kaggle/input'):
    for filename in filenames:
        print(os.path.join(dirname, filename))

# You can write up to 20GB to the current directory (/kaggle/working/) that gets preserved as output when you create a version using "Save & Run All" 
# You can also write temporary files to /kaggle/temp/, but they won't be saved outside of the current session


"""
capstone_ai_agent.py

A single-file capstone-ready project core:
 - InMemorySessionService (session + achievements)
 - MockLLMAgent (summaries, explanations, quiz generation)
 - ContentSearchAgent (resource finder - mock / extendable to web.run)
 - FlashcardAgent
 - StudyCoordinator (orchestrates agents, parallel tasks)
 - KaggleCoach (baseline pipeline, optional Kaggle API integration)
 - NotebookGenerator (writes reproducible .ipynb)
 - ReportGenerator (creates markdown capstone report skeleton)
 - PPTGenerator (creates simple PPTX slides using python-pptx if available)

Usage:
  1. Install dependencies:
     pip install pandas numpy scikit-learn joblib nbformat python-pptx kaggle

     - If you won't use PPT generation, python-pptx is optional.
     - If you won't use Kaggle API, 'kaggle' library is optional; KaggleCoach will run locally without submitting.

  2. Put your datasets in a working folder or use KaggleAgent.download_competition (requires kaggle.json).

  3. Run the demo at the bottom of the file or import the classes into your own scripts.

Notes on plagiarism and academic integrity:
  - This code is original and meant to be authored by you for a capstone.
  - When you create notebooks or write the report, ensure explanations, analyses,
    and textual writeups are written in your own words. Cite external sources where used.
"""

import os
import time
import uuid
import json
import shutil
from dataclasses import dataclass, asdict
from typing import List, Dict, Any, Optional, Tuple
import joblib
import nbformat
from nbformat.v4 import new_notebook, new_code_cell, new_markdown_cell

import numpy as np
import pandas as pd

from sklearn.model_selection import train_test_split, cross_val_score, StratifiedKFold, KFold
from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor
from sklearn.linear_model import LogisticRegression
from sklearn.metrics import accuracy_score, roc_auc_score, mean_squared_error, r2_score

# Try to import KaggleApi (optional). Kaggle functionality is wrapped so code works without it.
try:
    from kaggle.api.kaggle_api_extended import KaggleApi
except Exception:
    KaggleApi = None

# Try to import python-pptx for slide generation (optional)
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None

# ----------------------------
# Data classes
# ----------------------------
@dataclass
class Resource:
    title: str
    summary: str
    relevance: float = 0.5
    url: str = ""

@dataclass
class Flashcard:
    question: str
    answer: str
    ease: float = 2.5

@dataclass
class Achievement:
    title: str
    date: float
    description: str = ""
    source_url: str = ""

# ----------------------------
# Session service
# ----------------------------
class InMemorySessionService:
    """
    Simple session store for demo/capstone use.
    Stores per-session state and achievements.
    """
    def __init__(self):
        self.sessions: Dict[str, Dict[str, Any]] = {}

    def create_session(self, user_id: str, metadata: Optional[Dict[str, Any]] = None) -> str:
        sid = str(uuid.uuid4())
        self.sessions[sid] = {
            "user_id": user_id,
            "created_at": time.time(),
            "metadata": metadata or {},
            "state": {
                "topics": [],
                "weekly_hours": 0,
                "last_plan": None,
                "achievements": []
            }
        }
        return sid

    def get(self, session_id: str) -> Dict[str, Any]:
        return self.sessions.get(session_id, {})

    def set_state(self, session_id: str, key: str, value: Any):
        if session_id not in self.sessions:
            raise KeyError("Session not found")
        self.sessions[session_id]["state"][key] = value

    def get_state(self, session_id: str, key: str, default=None):
        return self.sessions.get(session_id, {}).get("state", {}).get(key, default)

    def add_achievement(self, session_id: str, achievement: Achievement):
        if session_id not in self.sessions:
            raise KeyError("Session not found")
        achs = self.sessions[session_id]["state"].get("achievements", [])
        achs.append(achievement)
        self.sessions[session_id]["state"]["achievements"] = achs

    def get_achievements(self, session_id: str) -> List[Achievement]:
        return self.get_state(session_id, "achievements", [])

# ----------------------------
# Mock LLM agent
# ----------------------------
class MockLLMAgent:
    """
    A simple, deterministic mock LLM to generate study text artifacts.
    For the capstone you can swap this out with actual LLM calls.
    """
    def summarize(self, text: str, max_length: int = 120) -> str:
        if not text:
            return ""
        s = text.strip().replace("\n", " ")
        return (s[:max_length].rstrip() + ("..." if len(s) > max_length else ""))

    def explain(self, topic: str, level: str = 'intro') -> str:
        return f"({level}) Explanation for {topic}: concise, learner-focused summary of main ideas."

    def generate_quiz(self, topic: str, n_questions: int = 3) -> List[Dict[str, str]]:
        qs = []
        for i in range(1, n_questions + 1):
            qs.append({"q": f"What is core idea {i} in {topic}?", "a": f"Short model answer {i} on {topic}."})
        return qs

    def answer_query(self, topic: str, query_type: str) -> str:
        if query_type == "prerequisites":
            return f"Prerequisites: basic Python, math foundations (linear algebra, probability), and ML basics for {topic}."
        if query_type == "key_concepts":
            return f"Key concepts: definitions, key algorithms, evaluation, and common pitfalls for {topic}."
        if query_type == "practice_problem":
            return f"Practice: implement a small project related to {topic}, solve targeted problems, and write a short writeup."
        return f"Short guidance on {topic} for {query_type}."

# ----------------------------
# Content search agent (mock)
# ----------------------------
class ContentSearchAgent:
    """
    Returns a small list of Resources for a topic. Replace or extend this with web scraping or API integration.
    """
    def run(self, topic: str) -> List[Resource]:
        # Simulated latency and content
        time.sleep(0.2)
        return [
            Resource(title=f"{topic} - Intro Article", summary=f"An introductory article about {topic}.", relevance=0.9, url=""),
            Resource(title=f"{topic} - Video Tutorial", summary=f"A concise tutorial video for {topic}.", relevance=0.85, url="")
        ]

# ----------------------------
# Flashcard agent
# ----------------------------
class FlashcardAgent:
    def __init__(self, llm: MockLLMAgent):
        self.llm = llm

    def run(self, topic: str, n_cards: int = 4) -> List[Flashcard]:
        cards = []
        for i in range(n_cards):
            q = f"What is a core idea {i+1} in {topic}?"
            a = self.llm.explain(topic, level='short')
            cards.append(Flashcard(question=q, answer=a))
        return cards

# ----------------------------
# Study coordinator
# ----------------------------
import concurrent.futures

class StudyCoordinator:
    """
    Coordinates content search, flashcard generation and LLM informative queries.
    Saves the 'plan' into session state.
    """
    def __init__(self, session_service: InMemorySessionService):
        self.llm = MockLLMAgent()
        self.search_agent = ContentSearchAgent()
        self.flashcard_agent = FlashcardAgent(self.llm)
        self.sessions = session_service

    def create_or_get_session(self, user_id: str, metadata: Optional[Dict[str, Any]] = None) -> str:
        # For simplicity: always create fresh session
        return self.sessions.create_session(user_id, metadata)

    def plan_for(self, session_id: str, topics: List[str], weekly_hours: int = 5) -> Dict[str, Any]:
        # store basics
        self.sessions.set_state(session_id, 'topics', topics)
        self.sessions.set_state(session_id, 'weekly_hours', weekly_hours)

        results = {"resources": {}, "flashcards": {}, "quizzes": {}, "informative_queries": {}}
        with concurrent.futures.ThreadPoolExecutor(max_workers=6) as ex:
            # schedule searches and flashcards
            search_futs = {ex.submit(self.search_agent.run, t): t for t in topics}
            flash_futs = {ex.submit(self.flashcard_agent.run, t, 4): t for t in topics}

            # collect search results
            for fut in concurrent.futures.as_completed(list(search_futs.keys())):
                t = search_futs[fut]
                try:
                    results['resources'][t] = fut.result()
                except Exception:
                    results['resources'][t] = []

            # collect flashcards
            for fut in concurrent.futures.as_completed(list(flash_futs.keys())):
                t = flash_futs[fut]
                try:
                    results['flashcards'][t] = fut.result()
                except Exception:
                    results['flashcards'][t] = []

            # sequentially ask LLM for quizzes and informative queries
            for t in topics:
                results['quizzes'][t] = self.llm.generate_quiz(t, n_questions=3)
                query_types = ["prerequisites", "key_concepts", "practice_problem"]
                iq = {}
                for qt in query_types:
                    iq[qt] = self.llm.answer_query(t, qt)
                results['informative_queries'][t] = iq

        # incorporate achievements: if Kaggle-style in achievements, append special flashcard
        achs = self.sessions.get_achievements(session_id)
        kaggle_like = any(('kaggle' in (a.title.lower() if isinstance(a, Achievement) else str(a.get('title','')).lower()))
                          for a in achs)
        if kaggle_like:
            for t in topics:
                extra = Flashcard(
                    question="How can competition experience be applied here?",
                    answer="Turn competition learnings into small reproducible experiments and document them.",
                    ease=3.0
                )
                results['flashcards'].setdefault(t, []).append(extra)

        self.sessions.set_state(session_id, 'last_plan', results)
        return results

# ----------------------------
# KaggleCoach - baseline pipeline
# ----------------------------
def ensure_dir(path: str):
    os.makedirs(path, exist_ok=True)
    return path

def safe_read_csv(path: str) -> pd.DataFrame:
    return pd.read_csv(path)

def is_classification_target(y: pd.Series) -> bool:
    return (y.dtype.name in ("int64", "int32") or y.dtype == object) and y.nunique() <= 30

class KaggleAgent:
    """
    Optional Kaggle API wrapper (best-effort). Requires kaggle package and kaggle.json.
    """
    def __init__(self, work_dir: str = "work"):
        self.work_dir = ensure_dir(work_dir)
        self.api = None
        if KaggleApi:
            try:
                self.api = KaggleApi()
                self.api.authenticate()
            except Exception:
                self.api = None

    def download_competition(self, comp_slug: str) -> str:
        if not self.api:
            raise RuntimeError("Kaggle API not configured/authenticated.")
        dest = ensure_dir(os.path.join(self.work_dir, comp_slug, "data"))
        print(f"[KaggleAgent] Downloading competition {comp_slug} to {dest}")
        self.api.competition_download_files(comp_slug, path=dest, quiet=False)
        # unzip zip files
        for f in os.listdir(dest):
            if f.endswith(".zip"):
                import zipfile
                with zipfile.ZipFile(os.path.join(dest, f), 'r') as zf:
                    zf.extractall(dest)
        return dest

    def submit(self, comp: str, file_path: str, message: str = "auto submission"):
        if not self.api:
            raise RuntimeError("Kaggle API not configured/authenticated.")
        print(f"[KaggleAgent] Submitting {file_path} to {comp}")
        self.api.competition_submit(file_path, message, comp)

    def get_leaderboard(self, comp: str):
        if not self.api:
            return None
        try:
            return self.api.competition_leaderboard_view(comp, page=1)
        except Exception:
            return None

class BaselinePipeline:
    """
    Simple, reproducible baseline pipeline:
      - basic missing-value handling, simple factorize for categoricals
      - baseline model (RandomForest)
      - cross-validation
      - saves model and submission if test provided
    """
    def __init__(self, work_dir: str = "work"):
        self.work_dir = ensure_dir(work_dir)

    def _artifact_path(self, comp: str, name: str) -> str:
        d = ensure_dir(os.path.join(self.work_dir, comp, "artifacts"))
        return os.path.join(d, name)

    def run(self, comp: str, train_csv: str, test_csv: Optional[str], target_col: str,
            random_state: int = 42) -> Dict[str, Any]:
        comp_dir = ensure_dir(os.path.join(self.work_dir, comp))
        print(f"[BaselinePipeline] Loading train: {train_csv}")
        train = safe_read_csv(train_csv)
        test = safe_read_csv(test_csv) if test_csv and os.path.exists(test_csv) else None

        if target_col not in train.columns:
            raise KeyError(f"Target column '{target_col}' not found in train data.")

        X = train.drop(columns=[target_col])
        y = train[target_col]

        # drop columns with > 80% missing
        miss_frac = X.isna().mean()
        drop_cols = miss_frac[miss_frac > 0.80].index.tolist()
        X = X.drop(columns=drop_cols)
        if test is not None:
            test = test.drop(columns=[c for c in drop_cols if c in test.columns], errors='ignore')

        # simple numeric imputation
        num_cols = X.select_dtypes(include=[np.number]).columns.tolist()
        for c in num_cols:
            m = X[c].median()
            X[c] = X[c].fillna(m)
            if test is not None and c in test.columns:
                test[c] = test[c].fillna(m)

        # simple categorical factorize
        cat_cols = X.select_dtypes(include=['object', 'category']).columns.tolist()
        cat_maps = {}
        for c in cat_cols:
            X[c], cats = pd.factorize(X[c].astype(str))
            cat_maps[c] = list(cats)
            if test is not None and c in test.columns:
                mapping = {v:i for i,v in enumerate(cats)}
                test[c] = test[c].map(mapping).fillna(-1).astype(int)

        task_type = "classification" if is_classification_target(y) else "regression"

        if task_type == "classification":
            model = RandomForestClassifier(n_estimators=200, random_state=random_state, n_jobs=-1)
            cv = StratifiedKFold(n_splits=5, shuffle=True, random_state=random_state)
            scoring = "accuracy"
        else:
            model = RandomForestRegressor(n_estimators=200, random_state=random_state, n_jobs=-1)
            cv = KFold(n_splits=5, shuffle=True, random_state=random_state)
            scoring = "neg_root_mean_squared_error"

        # cross-validate
        try:
            scores = cross_val_score(model, X, y, cv=cv, scoring=scoring, n_jobs=-1)
            cv_scores = scores.tolist()
        except Exception:
            cv_scores = []

        # fit on full data
        model.fit(X, y)
        model_path = self._artifact_path(comp, "baseline_model.joblib")
        joblib.dump(model, model_path)

        # create submission if test present
        submission_path = None
        if test is not None:
            # ensure we only use columns present in X
            common_cols = [c for c in X.columns if c in test.columns]
            preds = model.predict(test[common_cols])
            # if classification and predict_proba available, use for binary
            if task_type == "classification" and hasattr(model, "predict_proba") and preds is not None:
                try:
                    probs = model.predict_proba(test[common_cols])
                    # choose probability for class '1' if binary
                    if probs.shape[1] >= 2:
                        preds = probs[:, 1]
                except Exception:
                    pass
            # build a simple submission: use first id-like column if present
            id_col = None
            for cand in ["Id", "id", "ID", "PassengerId"]:
                if cand in test.columns:
                    id_col = cand
                    break
            df_sub = pd.DataFrame({"prediction": preds})
            if id_col:
                df_sub.insert(0, id_col, test[id_col].values)
            submission_path = os.path.join(comp_dir, "submission.csv")
            df_sub.to_csv(submission_path, index=False)
            print(f"[BaselinePipeline] submission written to {submission_path}")

        meta = {
            "n_train": len(train),
            "n_features": X.shape[1],
            "drop_columns": drop_cols,
            "cat_maps_summary": {k: len(v) for k,v in cat_maps.items()},
            "cv_scores": cv_scores,
            "scoring": scoring,
            "task_type": task_type,
            "model_path": model_path
        }

        return {"meta": meta, "submission": submission_path, "model_path": model_path}

# ----------------------------
# Notebook generator
# ----------------------------
class NotebookGenerator:
    """
    Create a short, reproducible notebook template (.ipynb) that documents the baseline pipeline steps.
    """
    def __init__(self, work_dir: str = "work"):
        self.work_dir = ensure_dir(work_dir)

    def generate(self, comp: str, train_csv: str, test_csv: Optional[str], target_col: str) -> str:
        nb = new_notebook()
        intro = f"# Auto-generated Notebook for {comp}\nThis notebook documents the baseline pipeline run.\n\nFiles: `{os.path.basename(train_csv)}`"
        if test_csv:
            intro += f", `{os.path.basename(test_csv)}`"
        nb.cells.append(new_markdown_cell(intro))

        nb.cells.append(new_code_cell(
            "import pandas as pd\nimport numpy as np\nfrom sklearn.model_selection import train_test_split\nfrom sklearn.ensemble import RandomForestClassifier\n"
        ))
        nb.cells.append(new_code_cell(f"train = pd.read_csv('{train_csv}')\ntrain.head()"))
        nb.cells.append(new_code_cell("train.describe(include='all').T"))
        nb.cells.append(new_code_cell(f"target = '{target_col}'\nX = train.drop(columns=[target])\ny = train[target]\nX_train, X_val, y_train, y_val = train_test_split(X, y, test_size=0.2, random_state=42)\nprint(X_train.shape)"))
        nb.cells.append(new_code_cell("# Add your feature engineering and modeling here\n"))

        outdir = ensure_dir(os.path.join(self.work_dir, comp, "notebooks"))
        fname = f"auto_notebook_{comp.replace(' ', '_')}.ipynb"
        path = os.path.join(outdir, fname)
        with open(path, "w", encoding="utf-8") as f:
            nbformat.write(nb, f)
        print(f"[NotebookGenerator] Notebook written to {path}")
        return path

# ----------------------------
# Report generator (markdown)
# ----------------------------
class ReportGenerator:
    """
    Create a markdown skeleton report for the capstone project. You should expand and fill with actual content.
    """
    def __init__(self, work_dir: str = "work"):
        self.work_dir = ensure_dir(work_dir)

    def generate_markdown(self, comp_title: str, authors: List[str], abstract: str, topics: List[str],
                          output_dir: Optional[str] = None) -> str:
        od = ensure_dir(output_dir or os.path.join(self.work_dir, "report"))
        fname = os.path.join(od, f"capstone_report_{comp_title.replace(' ', '_')}.md")
        md_lines = []
        md_lines.append(f"# {comp_title}\n")
        md_lines.append(f"**Authors:** {', '.join(authors)}\n")
        md_lines.append(f"**Date:** {time.strftime('%Y-%m-%d')}\n")
        md_lines.append("## Abstract\n")
        md_lines.append(abstract + "\n")
        md_lines.append("## Introduction\n")
        md_lines.append("Explain the problem, motivation and objectives.\n")
        md_lines.append("## System Architecture\n")
        md_lines.append("Describe modules: StudyCoordinator, Agents, SessionService, KaggleCoach.\n")
        md_lines.append("## Implementation Details\n")
        md_lines.append("Include code structure, libraries used, and important design choices.\n")
        md_lines.append("## Experiments and Results\n")
        md_lines.append("Document baseline experiments, CV scores, and submission results.\n")
        md_lines.append("## Discussion\n")
        md_lines.append("Limitations, ethics, future work.\n")
        md_lines.append("## Conclusion\n")
        md_lines.append("Summarize achievements and learning outcomes.\n")
        md_lines.append("## Appendix\n")
        md_lines.append("Include additional figures, hyperparameters, and dataset references.\n")

        with open(fname, "w", encoding="utf-8") as f:
            f.write("\n".join(md_lines))
        print(f"[ReportGenerator] Markdown report skeleton written to {fname}")
        return fname

# ----------------------------
# PPT generator (simple)
# ----------------------------
class PPTGenerator:
    """
    Very simple slide deck generator. Requires python-pptx.
    """
    def __init__(self, work_dir: str = "work"):
        self.work_dir = ensure_dir(work_dir)
        if Presentation is None:
            print("[PPTGenerator] python-pptx not available; PPT generation disabled.")

    def generate(self, title: str, subtitle: str, bullets: List[str], output_dir: Optional[str] = None) -> Optional[str]:
        if Presentation is None:
            print("[PPTGenerator] python-pptx not installed. Install with 'pip install python-pptx' to enable.")
            return None
        od = ensure_dir(output_dir or os.path.join(self.work_dir, "slides"))
        prs = Presentation()
        # title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        # bullet slide
        slide_b = prs.slides.add_slide(prs.slide_layouts[1])
        slide_b.shapes.title.text = "Key Points"
        body = slide_b.shapes.placeholders[1].text_frame
        for i, b in enumerate(bullets):
            if i == 0:
                body.text = b
            else:
                p = body.add_paragraph()
                p.text = b
        out = os.path.join(od, f"{title.replace(' ', '_')}.pptx")
        prs.save(out)
        print(f"[PPTGenerator] PPT saved to {out}")
        return out

# ----------------------------
# KaggleCoach wrapper tying pipeline + notebook + session
# ----------------------------
class KaggleCoach:
    def __init__(self, session_service: InMemorySessionService, work_dir: str = "work"):
        self.sessions = session_service
        self.kaggle_agent = KaggleAgent(work_dir=work_dir) if KaggleApi else None
        self.pipeline = BaselinePipeline(work_dir=work_dir)
        self.notebook_gen = NotebookGenerator(work_dir=work_dir)
        self.work_dir = work_dir

    def coach_run(self, session_id: str, comp_name: str, train_csv: str, test_csv: Optional[str],
                  target_col: str, submit: bool = False, swag_percentile: float = 20.0) -> Dict[str, Any]:
        # copy files to work/<comp>/data for reproducibility
        comp_data = ensure_dir(os.path.join(self.work_dir, comp_name, "data"))
        def copy_if_needed(fp: Optional[str]) -> Optional[str]:
            if not fp: return None
            dest = os.path.join(comp_data, os.path.basename(fp))
            if os.path.abspath(fp) != os.path.abspath(dest):
                shutil.copy(fp, dest)
            return dest

        tpath = copy_if_needed(train_csv)
        tepath = copy_if_needed(test_csv) if test_csv else None

        # run baseline
        res = self.pipeline.run(comp_name, tpath, tepath, target_col)

        # generate notebook
        nb = self.notebook_gen.generate(comp_name, tpath, tepath, target_col)
        res['notebook'] = nb

        # optionally submit using Kaggle API
        if submit and self.kaggle_agent:
            if res.get('submission'):
                try:
                    self.kaggle_agent.submit(comp_name, res['submission'])
                    res['submitted'] = True
                except Exception as e:
                    res['submitted'] = False
                    res['submit_error'] = str(e)
            else:
                res['submitted'] = False
                res['submit_error'] = "No submission file found."

        # best-effort: check leaderboard (mostly not available programmatically)
        rank_info = None
        if self.kaggle_agent:
            try:
                rank_info = self.kaggle_agent.get_leaderboard(comp_name)
            except Exception:
                rank_info = None
        res['public_leaderboard'] = rank_info

        # if hypothetical rank falls within swag_percentile, record an achievement (mock)
        # Note: this is conservative and mostly illustrative
        if isinstance(rank_info, dict) and 'rank' in rank_info:
            try:
                rank = int(rank_info['rank'])
                # Hypothetical check:
                if rank <= max(1, int(swag_percentile/100.0 * max(1000, rank))):
                    ach = Achievement(title=f"Kaggle Swag: rank {rank}", date=time.time(),
                                      description=f"Detected rank {rank} in {comp_name}")
                    self.sessions.add_achievement(session_id, ach)
                    res['achievement_recorded'] = True
            except Exception:
                pass

        self.sessions.set_state(session_id, "last_kaggle_run", res)
        return res

# ----------------------------
# Example demo & command-line-like usage
# ----------------------------
if __name__ == "__main__":
    print("Capstone AI Agent module loaded. Running minimal demo (no Kaggle API calls).")

    # create session service and coordinator
    ss = InMemorySessionService()
    sc = StudyCoordinator(ss)

    sid = sc.create_or_get_session(user_id="student_001", metadata={"preferred_format": "notebook+report"})
    print("Created session:", sid)

    # plan for 3 topics
    topics = ["Foundations of AI Agents", "Tools & Integrations", "Deployment & Observability"]
    plan = sc.plan_for(sid, topics, weekly_hours=8)
    print("Plan generated. Topics:", plan['resources'].keys())

    # Save an achievement to simulate Kaggle swag
    ach = Achievement(title="Kaggle - participation badge", date=time.time(), description="Simulated badge.")
    ss.add_achievement(sid, ach)

    # create kaggle coach (without using Kaggle API credentials here)
    coach = KaggleCoach(ss, work_dir="work")
    # To run a real baseline, set real CSV paths and target column. Example (uncomment and replace paths):
    # result = coach.coach_run(sid, "titanic_demo", "data/train.csv", "data/test.csv", target_col="Survived", submit=False)
    # print(json.dumps(result, indent=2))

    # generate report skeleton & ppt
    rg = ReportGenerator(work_dir="work")
    md = rg.generate_markdown(comp_title="AI Agent Study Planner - Capstone", authors=["Your Name"], 
                              abstract="This project builds an AI-agent-driven study planner and Kaggle coaching pipeline.",
                              topics=topics)
    print("Report skeleton written:", md)

    pptg = PPTGenerator(work_dir="work")
    ppt = pptg.generate(title="AI Agent Capstone", subtitle="Study Planner + Kaggle Coach", bullets=[
        "Multi-agent architecture (LLM, Search, Flashcards)",
        "Session persistence and achievements",
        "Baseline ML pipeline for competitions",
        "Notebook & report generation for reproducibility"
    ])
    print("PPT generated:", ppt)

    print("Demo finished. Inspect 'work/' directory for artifacts (notebooks, report, models).")


